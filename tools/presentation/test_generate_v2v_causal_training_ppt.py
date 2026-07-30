import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation

from tools.presentation.generate_v2v_causal_training_ppt import (
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    build_presentation,
    render_previews,
    validate_presentation,
)


class GeneratedDeckTest(unittest.TestCase):
    def test_generated_deck_has_expected_structure(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            path = build_presentation(Path(tmp_dir) / "deck.pptx")
            presentation = Presentation(path)

            self.assertEqual(len(presentation.slides), 15)
            self.assertEqual(presentation.slide_width, SLIDE_WIDTH)
            self.assertEqual(presentation.slide_height, SLIDE_HEIGHT)
            self.assertAlmostEqual(
                presentation.slide_width / presentation.slide_height,
                16 / 9,
                delta=0.01,
            )
            self.assertEqual(validate_presentation(path), [])

    def test_all_shapes_stay_inside_slide_bounds(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            path = build_presentation(Path(tmp_dir) / "deck.pptx")
            presentation = Presentation(path)
            for slide_number, slide in enumerate(presentation.slides, start=1):
                for shape in slide.shapes:
                    self.assertGreaterEqual(shape.left, 0, f"slide {slide_number}: {shape.name}")
                    self.assertGreaterEqual(shape.top, 0, f"slide {slide_number}: {shape.name}")
                    self.assertLessEqual(
                        shape.left + shape.width,
                        presentation.slide_width,
                        f"slide {slide_number}: {shape.name}",
                    )
                    self.assertLessEqual(
                        shape.top + shape.height,
                        presentation.slide_height,
                        f"slide {slide_number}: {shape.name}",
                    )

    def test_theme_colors_and_speaker_notes_are_embedded(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            path = build_presentation(Path(tmp_dir) / "deck.pptx")
            with zipfile.ZipFile(path) as archive:
                slide_xml = b"".join(
                    archive.read(name)
                    for name in archive.namelist()
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                )
                notes_xml = b"".join(
                    archive.read(name)
                    for name in archive.namelist()
                    if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")
                )

            for color in (b"14D9C4", b"FF9F43", b"8B7CFF", b"FF5C78"):
                self.assertIn(color, slide_xml)
            self.assertIn("开场先声明".encode(), notes_xml)

    def test_fallback_preview_renderer_outputs_all_slides_and_contact_sheet(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            root = Path(tmp_dir)
            path = build_presentation(root / "deck.pptx")
            preview_dir = root / "preview"
            contact_sheet = root / "contact-sheet.png"
            outputs = render_previews(path, preview_dir, contact_sheet)

            self.assertEqual(len(outputs), 15)
            self.assertTrue(all(output.exists() for output in outputs))
            self.assertTrue(contact_sheet.exists())


if __name__ == "__main__":
    unittest.main()
