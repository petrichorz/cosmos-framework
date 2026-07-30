#!/usr/bin/env python3
"""Generate the editable Cosmos3 V2V / causal-training presentation."""

from __future__ import annotations

import argparse
import re
import sys
from functools import lru_cache
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:
    from tools.presentation.v2v_causal_training_content import SLIDES, validate_content
except ModuleNotFoundError:
    from v2v_causal_training_content import SLIDES, validate_content


SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)

BG = "08131F"
PANEL = "102235"
PANEL_2 = "142B42"
GRID = "26445E"
TEXT = "F3F7FB"
MUTED = "9DB0C3"
CLEAN = "14D9C4"
NOISY = "FF9F43"
REASONER = "8B7CFF"
RISK = "FF5C78"
WHITE = "FFFFFF"
BLACK = "071018"
FONT = "Microsoft YaHei"
CJK_FONT_FILE = Path("/usr/share/fonts/google-droid-fonts/DroidSansFallback.ttf")
LATIN_FONT_FILE = Path("/usr/share/fonts/urw-base35/NimbusSans-Regular.otf")


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = TEXT,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.05,
    name: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    shape.text_frame.clear()
    shape.text_frame.margin_left = Inches(margin)
    shape.text_frame.margin_right = Inches(margin)
    shape.text_frame.margin_top = Inches(margin)
    shape.text_frame.margin_bottom = Inches(margin)
    shape.text_frame.vertical_anchor = valign
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = PANEL,
    line: str = GRID,
    radius: bool = True,
    transparency: int = 0,
    name: str | None = None,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: str = GRID, width: float = 1.5):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(min(x1, x2)),
        Inches(min(y1, y2)),
        Inches(max(abs(x2 - x1), 0.015)),
        Inches(max(abs(y2 - y1), 0.015)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_chevron(slide, x: float, y: float, w: float = 0.34, h: float = 0.38, color: str = MUTED):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_label(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    color: str,
    text_color: str = BLACK,
    size: float = 11,
):
    add_box(slide, x, y, w, 0.32, fill=color, line=color, radius=True)
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        0.32,
        size=size,
        color=text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0.01,
    )


def set_background(slide):
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(BG)
    for x, y, color, size in (
        (11.85, 0.05, CLEAN, 1.00),
        (12.30, 0.45, REASONER, 0.45),
    ):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.fill.transparency = 86
        shape.line.fill.background()


def add_header(slide, number: int, title: str, takeaway: str):
    add_text(slide, f"{number:02d}", 0.48, 0.35, 0.56, 0.38, size=12, color=CLEAN, bold=True)
    add_line(slide, 1.02, 0.43, 1.08, 0.67, CLEAN, 2)
    add_text(slide, title, 1.18, 0.28, 11.45, 0.58, size=25, color=TEXT, bold=True, margin=0)
    add_text(slide, takeaway, 0.58, 0.95, 12.10, 0.53, size=15, color=MUTED, margin=0)


def add_footer(slide, number: int, references: Iterable[str]):
    add_line(slide, 0.58, 7.03, 12.75, 7.045, GRID, 1)
    text = "  ·  ".join(references)
    add_text(slide, text, 0.58, 7.08, 11.78, 0.22, size=7.5, color=MUTED, valign=MSO_ANCHOR.TOP, margin=0)
    add_text(slide, f"{number:02d} / 15", 12.25, 7.06, 0.52, 0.22, size=8, color=CLEAN, align=PP_ALIGN.RIGHT)


def add_bullets(slide, bullets: list[str], x: float = 9.28, y: float = 1.75, w: float = 3.48, h: float = 4.95):
    add_box(slide, x, y, w, h, fill=PANEL, line=GRID)
    add_label(slide, "KEY POINTS", x + 0.20, y + 0.20, 1.18, color=CLEAN, size=9)
    cursor_y = y + 0.74
    item_h = (h - 1.0) / max(len(bullets), 1)
    colors = [CLEAN, NOISY, REASONER, RISK]
    for idx, bullet in enumerate(bullets):
        add_box(
            slide,
            x + 0.24,
            cursor_y + 0.08,
            0.12,
            0.12,
            fill=colors[idx % len(colors)],
            line=colors[idx % len(colors)],
            radius=False,
        )
        add_text(
            slide,
            bullet,
            x + 0.50,
            cursor_y - 0.05,
            w - 0.76,
            item_h,
            size=12.2,
            color=TEXT,
            valign=MSO_ANCHOR.TOP,
            margin=0,
        )
        cursor_y += item_h


def add_token(slide, text: str, x: float, y: float, w: float, color: str, *, h: float = 0.58, size: float = 12):
    add_box(slide, x, y, w, h, fill=color, line=color)
    add_text(slide, text, x, y, w, h, size=size, color=BLACK, bold=True, align=PP_ALIGN.CENTER, margin=0.01)


def add_attention_matrix(slide, x: float, y: float, n: int = 5, cell: float = 0.56):
    add_text(slide, "KEY  →  frame", x + 0.35, y - 0.38, n * cell, 0.28, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "QUERY", x - 0.62, y + 0.86, 0.46, 1.0, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    for row in range(n):
        add_text(slide, f"V{row}", x - 0.34, y + row * cell, 0.28, cell, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, f"V{row}", x + row * cell, y - 0.30, cell, 0.25, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        for col in range(n):
            visible = col <= row
            color = CLEAN if visible else PANEL_2
            add_box(
                slide,
                x + col * cell,
                y + row * cell,
                cell - 0.05,
                cell - 0.05,
                fill=color,
                line=color if visible else GRID,
                radius=False,
                transparency=0 if visible else 20,
            )
            add_text(
                slide,
                "✓" if visible else "×",
                x + col * cell,
                y + row * cell,
                cell - 0.05,
                cell - 0.05,
                size=12,
                color=BLACK if visible else MUTED,
                bold=visible,
                align=PP_ALIGN.CENTER,
            )


def draw_cover(slide, data):
    add_text(slide, "COSMOS3 EDGE · ALGORITHM REVIEW", 0.72, 0.62, 5.2, 0.36, size=12, color=CLEAN, bold=True)
    add_text(slide, data["title"], 0.72, 1.30, 8.3, 1.48, size=34, color=TEXT, bold=True, valign=MSO_ANCHOR.TOP)
    add_text(slide, data["takeaway"], 0.76, 3.00, 7.75, 0.70, size=17, color=MUTED, valign=MSO_ANCHOR.TOP)
    for i, (name, color, sub) in enumerate(
        (
            ("CONDITION", CLEAN, "哪些帧保持 clean"),
            ("ATTENTION", REASONER, "当前帧能看谁"),
            ("STRATEGY", NOISY, "历史处于何种噪声状态"),
        )
    ):
        y = 4.15 + i * 0.74
        add_label(slide, name, 0.78, y, 1.30, color=color, size=9)
        add_text(slide, sub, 2.25, y - 0.02, 3.7, 0.34, size=13, color=TEXT)
    add_box(slide, 9.08, 1.05, 3.40, 5.45, fill=PANEL, line=GRID)
    for t in range(5):
        frame_color = CLEAN if t == 0 else NOISY
        add_token(slide, f"V{t}", 9.62, 1.55 + t * 0.82, 0.78, frame_color, h=0.52)
        for p in range(3):
            add_box(
                slide,
                10.64 + p * 0.46,
                1.62 + t * 0.82,
                0.34,
                0.34,
                fill=frame_color,
                line=frame_color,
                transparency=12 * p,
            )
        if t < 4:
            add_line(slide, 10.02, 2.10 + t * 0.82, 10.02, 2.26 + t * 0.82, GRID, 1)
    add_text(slide, "V₀ clean  →  V₁…Vₜ generation", 9.45, 5.85, 2.72, 0.32, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "算法研发技术说明", 0.76, 6.65, 2.35, 0.30, size=11, color=TEXT)
    add_text(slide, "基线 ae9a530  ·  2026-07-30", 3.05, 6.65, 3.15, 0.30, size=11, color=MUTED)


def draw_capability(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    add_box(slide, 0.58, 1.75, 8.38, 4.95, fill=PANEL, line=GRID)
    add_label(slide, "CURRENT", 0.88, 2.02, 1.04, color=CLEAN, size=9)
    add_text(slide, "普通 V2V / I2V SFT", 0.88, 2.50, 3.6, 0.42, size=20, bold=True)
    add_token(slide, "clean V₀…Vₖ₋₁", 0.92, 3.18, 2.24, CLEAN)
    add_chevron(slide, 3.31, 3.28, color=MUTED)
    add_token(slide, "joint denoise", 3.80, 3.18, 1.72, NOISY)
    add_chevron(slide, 5.67, 3.28, color=MUTED)
    add_token(slide, "flow loss", 6.16, 3.18, 1.42, REASONER)
    add_line(slide, 0.90, 4.08, 8.62, 4.095, GRID, 1)
    add_label(slide, "GAP", 0.88, 4.34, 0.78, color=RISK, size=9)
    add_text(slide, "严格逐帧 causal training", 0.88, 4.83, 4.1, 0.42, size=20, bold=True)
    add_token(slide, "(T,S) causal", 0.92, 5.53, 1.70, REASONER)
    add_text(slide, "+", 2.74, 5.54, 0.32, 0.50, size=24, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_token(slide, "clean-history KV", 3.16, 5.53, 2.05, RISK)
    add_text(slide, "尚未在 Ascend 基础路径闭环", 5.45, 5.54, 2.82, 0.52, size=13, color=RISK, bold=True)
    add_footer(slide, data["number"], data["references"])


def draw_three_layers(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    layers = [
        ("01", "CONDITIONING", "σ_eff = 0 ?", "决定帧是否 clean", CLEAN),
        ("02", "ATTENTION MASK", "qₜ → k₀…kₜ ?", "决定能看哪些帧", REASONER),
        ("03", "TRAINING STRATEGY", "σ shared / per-frame / memory ?", "决定历史噪声与记忆", NOISY),
    ]
    for idx, (num, label, equation, desc, color) in enumerate(layers):
        y = 1.80 + idx * 1.56
        add_box(slide, 0.62, y, 8.30, 1.24, fill=PANEL, line=color)
        add_text(slide, num, 0.88, y + 0.18, 0.62, 0.50, size=24, color=color, bold=True)
        add_text(slide, label, 1.66, y + 0.15, 2.50, 0.34, size=13, color=color, bold=True)
        add_text(slide, equation, 1.66, y + 0.53, 3.42, 0.40, size=17, color=TEXT, bold=True)
        add_text(slide, desc, 5.35, y + 0.32, 3.05, 0.52, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, data["number"], data["references"])


def draw_condition_timeline(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    rows = [
        ("T2V", 0, "[]"),
        ("I2V", 1, "[0]"),
        ("V2V", 3, "[0,1,2]"),
    ]
    for row, (mode, clean_count, indexes) in enumerate(rows):
        y = 1.92 + row * 1.47
        add_text(slide, mode, 0.68, y + 0.10, 0.72, 0.42, size=16, color=TEXT, bold=True)
        add_text(slide, indexes, 1.46, y + 0.10, 1.02, 0.42, size=11, color=MUTED)
        for t in range(6):
            color = CLEAN if t < clean_count else NOISY
            add_token(slide, f"V{t}", 2.62 + t * 0.94, y, 0.72, color)
            if t < 5:
                add_line(slide, 3.34 + t * 0.94, y + 0.28, 3.53 + t * 0.94, y + 0.295, GRID, 1)
        add_text(
            slide,
            "clean prefix" if clean_count else "all generation",
            2.64,
            y + 0.66,
            2.65,
            0.24,
            size=9,
            color=CLEAN if clean_count else NOISY,
        )
    add_text(slide, "索引作用于 VAE latent frame，而非原始 pixel frame", 2.62, 6.35, 5.55, 0.34, size=12, color=MUTED)
    add_footer(slide, data["number"], data["references"])


def draw_pipeline(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    stages = [
        ("01", "SFTDataset", "video · text · plan", CLEAN),
        ("02", "VAE / text", "clean x₀ · token ids", REASONER),
        ("03", "PACK", "mask · t · loss idx", NOISY),
        ("04", "NOISE", "x₀ + ε + σ → xₜ", NOISY),
        ("05", "DENOISE", "Cosmos3VFMNetwork", REASONER),
        ("06", "LOSS", "flow matching vₜ", CLEAN),
    ]
    for idx, (num, title, sub, color) in enumerate(stages):
        col = idx % 3
        row = idx // 3
        x = 0.62 + col * 2.78
        y = 1.78 + row * 2.14
        add_box(slide, x, y, 2.30, 1.46, fill=PANEL, line=color)
        add_label(slide, num, x + 0.18, y + 0.18, 0.44, color=color, size=9)
        add_text(slide, title, x + 0.18, y + 0.57, 1.92, 0.34, size=16, bold=True)
        add_text(slide, sub, x + 0.18, y + 0.95, 1.92, 0.30, size=10, color=MUTED)
        if col < 2:
            add_chevron(slide, x + 2.41, y + 0.53, color=GRID)
    add_line(slide, 8.00, 3.10, 8.16, 3.55, GRID, 1)
    add_chevron(slide, 7.90, 3.35, 0.28, 0.32, GRID)
    add_footer(slide, data["number"], data["references"])


def draw_two_tower(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    add_box(slide, 0.62, 1.82, 3.45, 4.70, fill=PANEL, line=REASONER)
    add_label(slide, "REASONER", 0.92, 2.12, 1.26, color=REASONER, size=9)
    add_text(slide, "Understanding split", 0.92, 2.62, 2.60, 0.40, size=19, bold=True)
    for idx, text in enumerate(["[BOS]", "text₀", "text₁", "…", "[SOG]"]):
        add_token(slide, text, 0.94 + idx * 0.58, 3.40, 0.48, REASONER, h=0.50, size=9)
    add_text(slide, "causal self-attention", 1.02, 4.22, 2.52, 0.34, size=13, color=REASONER, align=PP_ALIGN.CENTER)
    add_text(slide, "不读取 Generator 中的视觉 token", 0.94, 5.23, 2.82, 0.52, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_box(slide, 5.05, 1.82, 3.84, 4.70, fill=PANEL, line=NOISY)
    add_label(slide, "GENERATOR", 5.36, 2.12, 1.34, color=NOISY, size=9)
    add_text(slide, "Generation split", 5.36, 2.62, 2.72, 0.40, size=19, bold=True)
    for t in range(4):
        add_token(slide, f"V{t}", 5.40 + t * 0.78, 3.40, 0.60, CLEAN if t == 0 else NOISY, h=0.50)
    add_text(slide, "clean + noisy VAE latent", 5.48, 4.22, 2.92, 0.34, size=13, color=NOISY, align=PP_ALIGN.CENTER)
    add_text(slide, "读取文本条件并预测 flow velocity", 5.35, 5.23, 3.05, 0.52, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_chevron(slide, 4.33, 3.55, 0.48, 0.62, CLEAN)
    add_text(slide, "text condition", 4.12, 4.24, 0.92, 0.28, size=9, color=CLEAN, align=PP_ALIGN.CENTER)
    add_footer(slide, data["number"], data["references"])


def draw_token_layout(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    add_label(slide, "UNDERSTANDING / CAUSAL", 0.72, 1.78, 2.30, color=REASONER, size=9)
    x = 0.72
    for text in ("BOS", "txt₀", "txt₁", "…", "EOS", "SOG"):
        add_token(slide, text, x, 2.25, 0.78, REASONER, h=0.52, size=9)
        x += 0.90
    add_label(slide, "GENERATOR / FULL", 0.72, 3.17, 1.82, color=NOISY, size=9)
    for t in range(3):
        base_x = 0.72 + t * 2.53
        add_text(slide, f"V{t}", base_x, 3.72, 0.42, 0.40, size=14, color=CLEAN if t == 0 else NOISY, bold=True)
        for p in range(4):
            color = CLEAN if t == 0 else NOISY
            add_box(slide, base_x + 0.45 + p * 0.43, 3.72, 0.34, 0.34, fill=color, line=color, transparency=p * 5)
        add_text(slide, "H → W", base_x + 0.45, 4.14, 1.66, 0.24, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_line(slide, 0.90, 5.03, 8.42, 5.045, GRID, 1)
    add_text(slide, "展平顺序", 0.74, 5.28, 0.94, 0.34, size=12, color=MUTED)
    add_text(slide, "T", 1.80, 5.28, 0.36, 0.34, size=15, color=CLEAN, bold=True)
    add_chevron(slide, 2.12, 5.28, 0.24, 0.28, GRID)
    add_text(slide, "H_patch", 2.48, 5.28, 0.88, 0.34, size=15, color=NOISY, bold=True)
    add_chevron(slide, 3.43, 5.28, 0.24, 0.28, GRID)
    add_text(slide, "W_patch", 3.80, 5.28, 0.92, 0.34, size=15, color=REASONER, bold=True)
    add_text(slide, "时间排序 ≠ temporal causal mask", 0.74, 6.15, 4.35, 0.34, size=13, color=RISK, bold=True)
    add_footer(slide, data["number"], data["references"])


def draw_mask_mapping(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    rows = [
        ("TOKEN", ["V₀", "V₁", "V₂", "V₃", "V₄"], [CLEAN, NOISY, NOISY, NOISY, NOISY]),
        ("condition_mask", ["1", "0", "0", "0", "0"], [CLEAN, PANEL_2, PANEL_2, PANEL_2, PANEL_2]),
        ("σ_eff", ["0", "σ", "σ", "σ", "σ"], [CLEAN, NOISY, NOISY, NOISY, NOISY]),
        ("loss", ["—", "✓", "✓", "✓", "✓"], [PANEL_2, REASONER, REASONER, REASONER, REASONER]),
    ]
    for row, (label, values, colors) in enumerate(rows):
        y = 1.88 + row * 1.02
        add_text(slide, label, 0.68, y + 0.11, 1.58, 0.42, size=12, color=MUTED, bold=True)
        for col, (value, color) in enumerate(zip(values, colors)):
            add_token(slide, value, 2.38 + col * 1.17, y, 0.88, color, h=0.58, size=12)
    add_box(slide, 0.72, 6.05, 7.95, 0.48, fill=PANEL_2, line=GRID)
    add_text(
        slide,
        "σ_eff = σ × (1 − condition_mask)     ·     loss 仅索引 generation patch",
        0.88,
        6.10,
        7.60,
        0.36,
        size=13,
        color=TEXT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, data["number"], data["references"])


def draw_supertoken(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    add_label(slide, "VISION ONLY", 0.72, 1.76, 1.26, color=CLEAN, size=9)
    for t in range(4):
        x = 0.72 + t * 1.96
        add_box(slide, x, 2.25, 1.60, 1.05, fill=PANEL, line=CLEAN)
        add_text(slide, f"Supertoken {t}", x, 2.35, 1.60, 0.26, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        add_token(slide, f"V{t} · H×W", x + 0.24, 2.72, 1.12, CLEAN if t == 0 else NOISY, h=0.40, size=9)
    add_label(slide, "ACTION + VISION", 0.72, 3.75, 1.58, color=REASONER, size=9)
    for t in range(4):
        x = 0.72 + t * 1.96
        add_box(slide, x, 4.22, 1.60, 1.28, fill=PANEL, line=REASONER)
        add_text(slide, f"t = {t}", x, 4.30, 1.60, 0.25, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        add_token(slide, "A" + str(t), x + 0.16, 4.78, 0.46, REASONER, h=0.40, size=9)
        add_token(slide, "V" + str(t), x + 0.71, 4.78, 0.72, CLEAN if t == 0 else NOISY, h=0.40, size=9)
    add_text(slide, "S = action tokens + H_patch × W_patch", 0.74, 6.08, 4.82, 0.34, size=13, color=MUTED)
    add_footer(slide, data["number"], data["references"])


def draw_attention(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    add_box(slide, 0.62, 1.75, 5.02, 4.95, fill=PANEL, line=GRID)
    add_attention_matrix(slide, 1.62, 2.50, n=5, cell=0.62)
    add_label(slide, "is_causal = ( True , False )", 1.20, 5.95, 3.84, color=CLEAN, size=11)
    add_box(slide, 5.92, 1.75, 3.02, 4.95, fill=PANEL, line=GRID)
    add_text(slide, "( T , S )", 6.34, 2.13, 2.16, 0.64, size=30, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "T 维", 6.35, 3.10, 0.70, 0.38, size=15, color=CLEAN, bold=True)
    add_text(slide, "CAUSAL", 7.16, 3.10, 1.10, 0.38, size=15, color=CLEAN, bold=True)
    add_text(slide, "S 维", 6.35, 3.78, 0.70, 0.38, size=15, color=NOISY, bold=True)
    add_text(slide, "FULL", 7.16, 3.78, 1.10, 0.38, size=15, color=NOISY, bold=True)
    add_line(slide, 6.30, 4.52, 8.58, 4.535, GRID, 1)
    add_label(slide, "REQUIRES", 6.38, 4.83, 1.08, color=RISK, size=9)
    add_text(slide, "three_way attention", 6.38, 5.32, 1.98, 0.52, size=14, color=RISK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, data["number"], data["references"])


def draw_strategy_table(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    rows, cols = 4, 5
    x, y, w, h = 0.60, 1.78, 12.15, 4.82
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    widths = [1.82, 1.78, 2.30, 2.30, 3.95]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    headers = ["strategy", "σ shape", "显式条件帧", "clean history", "当前语义"]
    matrix = [
        ["none", "[B,1]", "保留", "无", "标准联合去噪"],
        ["teacher_forcing", "[B,1]", "保留", "基础类无", "配置接口 + 空 hook"],
        ["diffusion_forcing", "[B,T]", "保留", "无", "逐 latent frame 独立 σ"],
    ]
    row_colors = [CLEAN, REASONER, NOISY]
    for col, value in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(PANEL_2)
        cell.text = value
    for row, values in enumerate(matrix, start=1):
        for col, value in enumerate(values):
            cell = table.cell(row, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(PANEL if row % 2 else PANEL_2)
            cell.text = value
            if col == 0:
                cell.fill.fore_color.rgb = rgb(row_colors[row - 1])
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(11.2 if row else 10.5)
                    run.font.bold = bool(row == 0 or col == 0)
                    run.font.color.rgb = rgb(BLACK if col == 0 and row else TEXT)
    add_text(
        slide,
        "注：teacher_forcing_dcm 目前仅作为配置接受值；本页不把它推断为已闭环能力。",
        0.68,
        6.68,
        11.9,
        0.25,
        size=10,
        color=MUTED,
    )
    add_footer(slide, data["number"], data["references"])


def draw_teacher_gap(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    add_box(slide, 0.62, 1.77, 4.04, 4.92, fill=PANEL, line=CLEAN)
    add_label(slide, "IDEAL", 0.90, 2.03, 0.78, color=CLEAN, size=9)
    add_text(slide, "预测 noisy V₂", 0.90, 2.52, 2.30, 0.40, size=19, bold=True)
    add_token(slide, "clean V₀", 0.92, 3.25, 1.04, CLEAN)
    add_token(slide, "clean V₁", 2.12, 3.25, 1.04, CLEAN)
    add_token(slide, "noisy V₂", 3.32, 3.25, 1.04, NOISY)
    add_text(slide, "clean-history forward / KV memory", 0.92, 4.18, 3.40, 0.42, size=13, color=CLEAN, align=PP_ALIGN.CENTER)
    add_box(slide, 4.92, 1.77, 4.04, 4.92, fill=PANEL, line=RISK)
    add_label(slide, "CURRENT BASE", 5.20, 2.03, 1.40, color=RISK, size=9)
    add_text(slide, "预测 noisy V₂", 5.20, 2.52, 2.30, 0.40, size=19, bold=True)
    add_token(slide, "clean V₀", 5.22, 3.25, 1.04, CLEAN)
    add_token(slide, "noisy V₁", 6.42, 3.25, 1.04, NOISY)
    add_token(slide, "noisy V₂", 7.62, 3.25, 1.04, NOISY)
    add_text(slide, "hook 原样返回  ·  memory=None", 5.22, 4.18, 3.38, 0.42, size=13, color=RISK, align=PP_ALIGN.CENTER)
    add_text(slide, "≠", 4.45, 3.30, 0.40, 0.52, size=28, color=RISK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, data["number"], data["references"])


def draw_combination_matrix(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    add_text(slide, "ATTENTION", 0.66, 1.77, 1.20, 0.30, size=10, color=MUTED, bold=True)
    add_text(slide, "FULL / FUTURE VISIBLE", 2.25, 1.77, 2.55, 0.30, size=11, color=REASONER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "TEMPORAL CAUSAL", 5.47, 1.77, 2.55, 0.30, size=11, color=CLEAN, bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ("SHARED σ", "联合去噪\n未来可见", "只看过去\n历史通常 noisy"),
        ("PER-FRAME σ", "逐帧噪声\n未来仍可见", "只看过去\n历史 σ 各异"),
    ]
    for row, (label, left, right) in enumerate(rows):
        y = 2.30 + row * 1.98
        add_text(slide, label, 0.66, y + 0.57, 1.26, 0.42, size=11, color=NOISY, bold=True)
        add_box(slide, 2.18, y, 2.82, 1.56, fill=PANEL, line=REASONER)
        add_text(slide, left, 2.38, y + 0.22, 2.42, 1.08, size=16, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_box(slide, 5.32, y, 2.82, 1.56, fill=PANEL, line=CLEAN)
        add_text(slide, right, 5.52, y + 0.22, 2.42, 1.08, size=16, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_label(slide, "MISSING CELL", 2.18, 6.23, 1.40, color=RISK, size=9)
    add_text(slide, "causal attention + clean-history memory", 3.78, 6.18, 4.10, 0.42, size=14, color=RISK, bold=True)
    add_footer(slide, data["number"], data["references"])


def draw_roadmap(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    add_bullets(slide, data["bullets"])
    stages = [
        ("NOW", "普通 V2V / I2V SFT", "two_way\nnon-causal\nstrategy=none", CLEAN),
        ("A", "NPU temporal attention", "(T,S) mask\nvarlen/padded\nforward + backward", REASONER),
        ("B", "Clean-history memory", "clean forward\nKV visibility\n显存与梯度策略", NOISY),
    ]
    for idx, (tag, title, body, color) in enumerate(stages):
        x = 0.62 + idx * 2.78
        add_box(slide, x, 1.88, 2.36, 3.95, fill=PANEL, line=color)
        add_label(slide, tag, x + 0.22, 2.12, 0.62, color=color, size=9)
        add_text(slide, title, x + 0.22, 2.72, 1.92, 0.74, size=17, color=TEXT, bold=True, valign=MSO_ANCHOR.TOP)
        add_line(slide, x + 0.22, 3.63, x + 2.12, 3.645, GRID, 1)
        add_text(slide, body, x + 0.22, 3.92, 1.92, 1.34, size=13, color=MUTED, valign=MSO_ANCHOR.TOP)
        if idx < 2:
            add_chevron(slide, x + 2.45, 3.46, 0.25, 0.40, GRID)
    add_box(slide, 0.62, 6.12, 7.92, 0.45, fill=PANEL_2, line=GRID)
    add_text(
        slide,
        "严格 V2V causal training = 阶段 A（可见性）+ 阶段 B（clean history）",
        0.80,
        6.15,
        7.55,
        0.36,
        size=13,
        color=TEXT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, data["number"], data["references"])


def draw_summary(slide, data):
    add_header(slide, data["number"], data["title"], data["takeaway"])
    questions = [
        ("01", "帧是否 clean？", "condition_frame_indexes_vision", CLEAN),
        ("02", "能看哪些帧？", "full vs temporal causal", REASONER),
        ("03", "历史是什么状态？", "shared σ / per-frame σ / clean KV", NOISY),
    ]
    for idx, (num, question, answer, color) in enumerate(questions):
        x = 0.64 + idx * 4.15
        add_box(slide, x, 1.92, 3.72, 2.22, fill=PANEL, line=color)
        add_text(slide, num, x + 0.22, 2.16, 0.54, 0.46, size=23, color=color, bold=True)
        add_text(slide, question, x + 0.22, 2.79, 3.28, 0.42, size=18, color=TEXT, bold=True)
        add_text(slide, answer, x + 0.22, 3.37, 3.28, 0.36, size=11, color=MUTED)
    add_box(slide, 0.64, 4.55, 12.02, 1.58, fill=PANEL_2, line=GRID)
    for idx, bullet in enumerate(data["bullets"]):
        x = 0.90 + (idx % 2) * 5.92
        y = 4.83 + (idx // 2) * 0.58
        add_box(slide, x, y + 0.07, 0.12, 0.12, fill=[CLEAN, REASONER, NOISY, RISK][idx], line=[CLEAN, REASONER, NOISY, RISK][idx], radius=False)
        add_text(slide, bullet, x + 0.26, y, 5.30, 0.34, size=12, color=TEXT)
    add_text(slide, "DISCUSSION", 0.66, 6.48, 1.08, 0.28, size=9, color=CLEAN, bold=True)
    add_text(slide, "目标是联合去噪式 V2V，还是严格逐帧 causal generation？", 1.72, 6.42, 7.48, 0.36, size=13, color=MUTED)
    add_footer(slide, data["number"], data["references"])


DRAWERS = {
    "cover": draw_cover,
    "capability": draw_capability,
    "three_layers": draw_three_layers,
    "condition_timeline": draw_condition_timeline,
    "pipeline": draw_pipeline,
    "two_tower": draw_two_tower,
    "token_layout": draw_token_layout,
    "mask_mapping": draw_mask_mapping,
    "supertoken": draw_supertoken,
    "attention_matrix": draw_attention,
    "strategy_table": draw_strategy_table,
    "teacher_forcing_gap": draw_teacher_gap,
    "combination_matrix": draw_combination_matrix,
    "roadmap": draw_roadmap,
    "summary": draw_summary,
}


def add_speaker_notes(slide, notes: str):
    text_frame = slide.notes_slide.notes_text_frame
    text_frame.text = notes
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT


def build_presentation(output_path: Path) -> Path:
    """Build and save the 15-slide editable PowerPoint deck."""

    content_errors = validate_content(SLIDES)
    if content_errors:
        raise ValueError("content validation failed: " + "; ".join(content_errors))

    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    blank_layout = presentation.slide_layouts[6]

    for data in SLIDES:
        slide = presentation.slides.add_slide(blank_layout)
        set_background(slide)
        DRAWERS[data["visual"]](slide, data)
        add_speaker_notes(slide, data["notes"])

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path


def _shape_fill_hex(shape) -> str | None:
    try:
        value = shape.fill.fore_color.rgb
        return str(value) if value is not None else None
    except (AttributeError, TypeError, ValueError):
        return None


def _shape_line_hex(shape) -> str | None:
    try:
        value = shape.line.color.rgb
        return str(value) if value is not None else None
    except (AttributeError, TypeError, ValueError):
        return None


@lru_cache(maxsize=64)
def _font(size: int, *, latin: bool = False) -> ImageFont.FreeTypeFont:
    font_file = LATIN_FONT_FILE if latin else CJK_FONT_FILE
    return ImageFont.truetype(str(font_file), max(8, size))


def _font_for_char(char: str, size: int) -> ImageFont.FreeTypeFont:
    codepoint = ord(char) if char else 0
    is_cjk = (
        0x2E80 <= codepoint <= 0x9FFF
        or 0xF900 <= codepoint <= 0xFAFF
        or 0xFF00 <= codepoint <= 0xFFEF
    )
    return _font(size, latin=not is_cjk)


def _mixed_text_length(draw: ImageDraw.ImageDraw, text: str, size: int) -> float:
    return sum(draw.textlength(char, font=_font_for_char(char, size)) for char in text)


def _draw_mixed_text(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, size: int, fill: str):
    cursor_x, cursor_y = xy
    for char in text:
        font = _font_for_char(char, size)
        draw.text((cursor_x, cursor_y), char, font=font, fill=fill)
        cursor_x += round(draw.textlength(char, font=font))


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font_size: int, max_width: int) -> list[str]:
    lines: list[str] = []
    for source_line in text.splitlines() or [""]:
        if not source_line:
            lines.append("")
            continue
        current = ""
        for char in source_line:
            candidate = current + char
            if current and _mixed_text_length(draw, candidate, font_size) > max_width:
                lines.append(current)
                current = char
            else:
                current = candidate
        if current:
            lines.append(current)
    return lines


def _draw_shape_text(draw: ImageDraw.ImageDraw, shape, box: tuple[int, int, int, int], scale: float):
    if not getattr(shape, "has_text_frame", False) or not shape.text:
        return
    x0, y0, x1, y1 = box
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.runs[0] if paragraph.runs else None
    point_size = float(run.font.size.pt) if run is not None and run.font.size else 14.0
    font_size = max(8, round(point_size * scale / 72))
    color = TEXT
    if run is not None:
        try:
            if run.font.color.rgb is not None:
                color = str(run.font.color.rgb)
        except (AttributeError, TypeError):
            pass
    max_width = max(8, x1 - x0 - round(0.10 * scale))
    lines = _wrap_text(draw, shape.text, font_size, max_width)
    line_height = max(10, round(font_size * 1.22))
    text_height = line_height * len(lines)
    align = paragraph.alignment
    valign = shape.text_frame.vertical_anchor
    if valign == MSO_ANCHOR.TOP:
        cursor_y = y0 + round(0.06 * scale)
    else:
        cursor_y = y0 + max(0, (y1 - y0 - text_height) // 2)
    for line in lines:
        line_width = _mixed_text_length(draw, line, font_size)
        if align == PP_ALIGN.CENTER:
            cursor_x = x0 + max(0, (x1 - x0 - round(line_width)) // 2)
        elif align == PP_ALIGN.RIGHT:
            cursor_x = x1 - round(line_width) - round(0.06 * scale)
        else:
            cursor_x = x0 + round(0.06 * scale)
        _draw_mixed_text(draw, (cursor_x, cursor_y), line, font_size, f"#{color}")
        cursor_y += line_height


def _draw_table(draw: ImageDraw.ImageDraw, shape, x0: int, y0: int, scale: float):
    table = shape.table
    cursor_y = y0
    for row in table.rows:
        row_h = round(row.height / 914400 * scale)
        cursor_x = x0
        for cell, column in zip(row.cells, table.columns):
            col_w = round(column.width / 914400 * scale)
            fill = PANEL
            try:
                if cell.fill.fore_color.rgb is not None:
                    fill = str(cell.fill.fore_color.rgb)
            except (AttributeError, TypeError):
                pass
            draw.rectangle((cursor_x, cursor_y, cursor_x + col_w, cursor_y + row_h), fill=f"#{fill}", outline=f"#{GRID}")
            run = cell.text_frame.paragraphs[0].runs[0] if cell.text_frame.paragraphs[0].runs else None
            point_size = float(run.font.size.pt) if run is not None and run.font.size else 11.0
            font_size = max(8, round(point_size * scale / 72))
            color = TEXT
            if run is not None:
                try:
                    if run.font.color.rgb is not None:
                        color = str(run.font.color.rgb)
                except (AttributeError, TypeError):
                    pass
            lines = _wrap_text(draw, cell.text, font_size, max(8, col_w - round(0.12 * scale)))
            line_h = max(10, round(font_size * 1.18))
            text_h = len(lines) * line_h
            text_y = cursor_y + max(0, (row_h - text_h) // 2)
            for line in lines:
                text_w = _mixed_text_length(draw, line, font_size)
                _draw_mixed_text(
                    draw,
                    (cursor_x + max(3, (col_w - round(text_w)) // 2), text_y),
                    line,
                    font_size,
                    f"#{color}",
                )
                text_y += line_h
            cursor_x += col_w
        cursor_y += row_h


def render_previews(
    pptx_path: Path,
    preview_dir: Path,
    contact_sheet_path: Path,
    *,
    pixels_per_inch: int = 120,
) -> list[Path]:
    """Render a deterministic Pillow preview from PPTX geometry when LibreOffice is unavailable."""

    presentation = Presentation(pptx_path)
    preview_dir.mkdir(parents=True, exist_ok=True)
    contact_sheet_path.parent.mkdir(parents=True, exist_ok=True)
    slide_width = round(presentation.slide_width / 914400 * pixels_per_inch)
    slide_height = round(presentation.slide_height / 914400 * pixels_per_inch)
    outputs: list[Path] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        image = Image.new("RGB", (slide_width, slide_height), f"#{BG}")
        draw = ImageDraw.Draw(image)
        for shape in slide.shapes:
            x0 = round(shape.left / 914400 * pixels_per_inch)
            y0 = round(shape.top / 914400 * pixels_per_inch)
            x1 = round((shape.left + shape.width) / 914400 * pixels_per_inch)
            y1 = round((shape.top + shape.height) / 914400 * pixels_per_inch)
            if shape.has_table:
                _draw_table(draw, shape, x0, y0, pixels_per_inch)
                continue
            fill = _shape_fill_hex(shape)
            line = _shape_line_hex(shape)
            if fill:
                if getattr(shape, "auto_shape_type", None) == MSO_SHAPE.OVAL:
                    draw.ellipse((x0, y0, x1, y1), fill=f"#{fill}", outline=f"#{line or fill}")
                elif getattr(shape, "auto_shape_type", None) == MSO_SHAPE.ROUNDED_RECTANGLE:
                    draw.rounded_rectangle(
                        (x0, y0, x1, y1),
                        radius=max(3, round(0.08 * pixels_per_inch)),
                        fill=f"#{fill}",
                        outline=f"#{line or fill}",
                    )
                elif getattr(shape, "auto_shape_type", None) == MSO_SHAPE.CHEVRON:
                    mid = (y0 + y1) // 2
                    points = [(x0, y0), (x1 - (y1 - y0) // 3, y0), (x1, mid), (x1 - (y1 - y0) // 3, y1), (x0, y1), (x0 + (y1 - y0) // 3, mid)]
                    draw.polygon(points, fill=f"#{fill}")
                else:
                    draw.rectangle((x0, y0, x1, y1), fill=f"#{fill}", outline=f"#{line or fill}")
            _draw_shape_text(draw, shape, (x0, y0, x1, y1), pixels_per_inch)
        output = preview_dir / f"slide-{slide_number:02d}.png"
        image.save(output, quality=95)
        outputs.append(output)

    thumb_w, thumb_h = 432, 243
    gutter, label_h = 22, 28
    sheet = Image.new("RGB", (3 * thumb_w + 4 * gutter, 5 * (thumb_h + label_h) + 6 * gutter), "#050D15")
    sheet_draw = ImageDraw.Draw(sheet)
    label_font = _font(15, latin=True)
    for idx, output in enumerate(outputs):
        row, col = divmod(idx, 3)
        x = gutter + col * (thumb_w + gutter)
        y = gutter + row * (thumb_h + label_h + gutter)
        thumb = Image.open(output).resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        sheet.paste(thumb, (x, y))
        sheet_draw.text((x, y + thumb_h + 4), f"{idx + 1:02d}", font=label_font, fill=f"#{CLEAN}")
    sheet.save(contact_sheet_path, quality=95)
    return outputs


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            parts.append(shape.text)
        if shape.has_table:
            for row in shape.table.rows:
                parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def validate_presentation(path: Path) -> list[str]:
    """Validate slide count, aspect ratio, bounds, references, notes, and placeholders."""

    errors: list[str] = []
    presentation = Presentation(path)
    if len(presentation.slides) != 15:
        errors.append(f"expected 15 slides, got {len(presentation.slides)}")
    if abs(presentation.slide_width / presentation.slide_height - 16 / 9) > 0.01:
        errors.append("slide aspect ratio is not 16:9")

    placeholder = re.compile(r"\b(?:TBD|TODO)\b|待补充|占位文本", re.IGNORECASE)
    for number, slide in enumerate(presentation.slides, start=1):
        text = _slide_text(slide)
        if number > 1 and f"{number:02d}" not in text:
            errors.append(f"slide {number} is missing page number")
        if "cosmos_framework/" not in text and number > 1:
            errors.append(f"slide {number} is missing code references")
        if placeholder.search(text):
            errors.append(f"slide {number} contains placeholder text")
        if not slide.notes_slide.notes_text_frame.text.strip():
            errors.append(f"slide {number} has no speaker notes")
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                errors.append(f"slide {number} shape {shape.name} starts outside bounds")
            if shape.left + shape.width > presentation.slide_width:
                errors.append(f"slide {number} shape {shape.name} exceeds right bound")
            if shape.top + shape.height > presentation.slide_height:
                errors.append(f"slide {number} shape {shape.name} exceeds bottom bound")
    return errors


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("docs/presentations/v2v_causal_training_zh.pptx"),
        help="Output PPTX path",
    )
    parser.add_argument("--check", type=Path, help="Validate an existing PPTX instead of generating")
    args = parser.parse_args(argv)

    path = args.check if args.check else build_presentation(args.output)
    errors = validate_presentation(path)
    if errors:
        print("Validation failed:", file=sys.stderr)
        for error in errors:
            print(f"- {error}", file=sys.stderr)
        return 1
    print(f"OK: {path} · 15 slides · 16:9 · content/notes/bounds validated")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
